# =========================================================
# creative_visual_tools.py — Advanced Branded AI Visual Builder (v3)
# =========================================================
import os
from datetime import datetime
import matplotlib.pyplot as plt
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

from AI_SYSTEM.AGENTS.CREATIVE_AGENT.tools.creative_tools import (
    load_sales_master, summarize_sales,
    load_returns_master, summarize_returns,
    _hex_to_rgb, choose_theme, load_profile
)
from AI_SYSTEM.RAG.rag_brain import UnifiedRAGBrain

# Optional image generator for theme backgrounds
try:
    from image_gen import text2im
    HAVE_AI_IMG = True
except Exception:
    HAVE_AI_IMG = False

BASE_PATH = r"C:\Users\Sandeep\Desktop\roborana_ai_system\RoboRana_AI_Data"
OUTPUT_DIR = os.path.join(BASE_PATH, "OUTPUTS")
LOGO_PATH = os.path.join(BASE_PATH, "AI_SYSTEM", "ASSETS", "logo.png")  # optional logo
os.makedirs(OUTPUT_DIR, exist_ok=True)

# ---------------------------------------------------------
# 🧱 Visual Components
# ---------------------------------------------------------
def _add_logo_and_watermark(slide, theme):
    """Add company logo and watermark on each slide."""
    try:
        if os.path.exists(LOGO_PATH):
            slide.shapes.add_picture(LOGO_PATH, Inches(8.5), Inches(0.2), height=Inches(0.7))
    except Exception:
        pass

    tx = slide.shapes.add_textbox(Inches(0.3), Inches(6.7), Inches(9), Inches(0.5))
    p = tx.text_frame.paragraphs[0]
    p.text = "© RoboRana AI — Automated Insight System"
    p.font.size = Pt(10)
    p.font.color.rgb = RGBColor(*_hex_to_rgb(theme.get("muted", "#999999")))
    p.alignment = PP_ALIGN.RIGHT

def generate_theme_background(theme_name: str, out_path: str):
    """AI-generate a theme-appropriate background."""
    if not HAVE_AI_IMG:
        return None

    prompts = {
        "floral": "pastel watercolor floral background with soft pink and purple tones",
        "modern": "minimal geometric gradient background with blue and grey tones, professional corporate style",
        "midnight": "dark abstract neon gradient background with deep blue highlights",
        "sunset": "warm glowing orange-yellow gradient with soft lighting and abstract texture",
    }

    prompt = prompts.get(theme_name.lower(), "minimal gradient background for presentation, modern style")
    try:
        text2im({"prompt": prompt, "size": "1024x768", "n": 1})
        # You can later modify this to actually save a generated image in your system
    except Exception as e:
        print(f"⚠️ Background generation failed: {e}")
    return out_path

# ---------------------------------------------------------
# 🎨 Slide Builders
# ---------------------------------------------------------
def _add_cover_slide(prs, title, subtitle, theme, bg_img=None):
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank slide
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(*_hex_to_rgb(theme.get("bg", "#FFFFFF")))

    # background AI image
    if bg_img and os.path.exists(bg_img):
        slide.shapes.add_picture(bg_img, 0, 0, width=prs.slide_width, height=prs.slide_height)

    # Title text
    tx = slide.shapes.add_textbox(Inches(0.8), Inches(1.8), Inches(8), Inches(2))
    tf = tx.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(44)
    p.font.bold = True
    p.font.color.rgb = RGBColor(*_hex_to_rgb(theme.get("primary", "#000000")))

    p2 = tf.add_paragraph()
    p2.text = subtitle
    p2.font.size = Pt(20)
    p2.font.color.rgb = RGBColor(*_hex_to_rgb(theme.get("accent", "#555555")))

    _add_logo_and_watermark(slide, theme)


def _add_kpi_slide(prs, heading, kpis: dict, theme):
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    slide.shapes.title.text = heading
    slide.shapes.title.text_frame.paragraphs[0].font.size = Pt(32)
    slide.shapes.title.text_frame.paragraphs[0].font.color.rgb = RGBColor(*_hex_to_rgb(theme.get("primary", "#000000")))

    left, top = Inches(0.7), Inches(1.8)
    for idx, (k, v) in enumerate(kpis.items()):
        y = top + Inches(idx * 0.9)
        tb = slide.shapes.add_textbox(left, y, Inches(8), Inches(0.8))
        p = tb.text_frame.paragraphs[0]
        p.text = f"{k}: {v}"
        p.font.size = Pt(24)
        p.font.color.rgb = RGBColor(*_hex_to_rgb(theme.get("text", "#000000")))

    _add_logo_and_watermark(slide, theme)


def _add_chart_slide(prs, heading, data: dict, theme):
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    slide.shapes.title.text = heading
    slide.shapes.title.text_frame.paragraphs[0].font.size = Pt(28)

    if not data:
        tb = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(8), Inches(3))
        tb.text = "No chart data available."
        return

    tmp_chart = os.path.join(OUTPUT_DIR, "TEMP", f"{heading.replace(' ', '_')}.png")
    os.makedirs(os.path.dirname(tmp_chart), exist_ok=True)

    plt.figure(figsize=(7, 4))
    plt.barh(list(data.keys())[:10], list(data.values())[:10], color=theme.get("primary", "#1F4E79"))
    plt.gca().invert_yaxis()
    plt.title(heading)
    plt.tight_layout()
    plt.savefig(tmp_chart, dpi=200)
    plt.close()

    slide.shapes.add_picture(tmp_chart, Inches(1), Inches(1.8), width=Inches(8))
    _add_logo_and_watermark(slide, theme)


def _add_ai_summary_slide(prs, theme, context: str = None):
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    slide.shapes.title.text = "AI Insights Summary"
    tf = slide.shapes.title.text_frame
    tf.paragraphs[0].font.size = Pt(28)

    tx = slide.shapes.add_textbox(Inches(0.8), Inches(1.8), Inches(8.5), Inches(4))
    body = tx.text_frame
    body.text = context or (
        "RoboRana AI suggests focusing on improving visibility and channel ROI. "
        "High-value products show steady growth, while returns remain stable."
    )
    body.paragraphs[0].font.size = Pt(18)
    body.paragraphs[0].font.color.rgb = RGBColor(*_hex_to_rgb(theme.get("text", "#000000")))
    _add_logo_and_watermark(slide, theme)

# ---------------------------------------------------------
# 🧱 Builder Functions
# ---------------------------------------------------------
def build_sales_ppt_v3(days=7, theme_key="modern", additional_context=None):
    profile = load_profile()
    theme_key, theme = choose_theme(profile, theme_key)

    df = load_sales_master()
    s = summarize_sales(df, days)
    prs = Presentation()
    rag = UnifiedRAGBrain()

    bg_path = os.path.join(OUTPUT_DIR, "TEMP", f"{theme_key}_bg.png")
    generate_theme_background(theme_key, bg_path)

    now = datetime.now().strftime("%d %b %Y, %H:%M")
    _add_cover_slide(prs, f"Sales Report — Last {days} Days", f"Generated {now}", theme, bg_path)
    _add_kpi_slide(prs, "Performance Overview", {
        "Orders": f"{s['total_orders']:,}",
        "Revenue": f"₹{s['total_value']:,.2f}",
        "AOV": f"₹{s['aov']:,.2f}"
    }, theme)
    if s["channels"]:
        _add_chart_slide(prs, "Top Channels by Revenue", s["channels"], theme)

    # RAG summary (AI-driven)
    try:
        insight = rag.query("Summarize last week sales insights for management")
    except Exception:
        insight = "AI Insight generation temporarily unavailable."

    _add_ai_summary_slide(prs, theme, insight)

    # Add Hive Mind context if provided
    if additional_context:
        try:
            slide = prs.slides.add_slide(prs.slide_layouts[5])
            slide.shapes.title.text = "Hive Mind Context Summary"
            tx = slide.shapes.add_textbox(Inches(0.8), Inches(1.8), Inches(8.5), Inches(4))
            body = tx.text_frame
            body.text = additional_context[:2000]
            body.paragraphs[0].font.size = Pt(16)
            body.paragraphs[0].font.color.rgb = RGBColor(*_hex_to_rgb(theme.get("text", "#000000")))
            _add_logo_and_watermark(slide, theme)
        except Exception as e:
            print(f"⚠️ Could not add context slide: {e}")

    out_path = os.path.join(OUTPUT_DIR, "PPT", f"Sales_Report_v3_{theme_key}_theme.pptx")
    os.makedirs(os.path.dirname(out_path), exist_ok=True)
    prs.save(out_path)
    return out_path


def build_returns_ppt_v3(days=7, theme_key="modern", additional_context=None):
    profile = load_profile()
    theme_key, theme = choose_theme(profile, theme_key)

    df = load_returns_master()
    r = summarize_returns(df, days)
    prs = Presentation()
    rag = UnifiedRAGBrain()

    bg_path = os.path.join(OUTPUT_DIR, "TEMP", f"{theme_key}_bg.png")
    generate_theme_background(theme_key, bg_path)

    now = datetime.now().strftime("%d %b %Y, %H:%M")
    _add_cover_slide(prs, f"Returns Report — Last {days} Days", f"Generated {now}", theme, bg_path)
    _add_kpi_slide(prs, "Key Metrics", {
        "Returns": f"{r['rows']:,}",
        "Return Value": f"₹{r['value']:,.2f}"
    }, theme)
    if r["channels"]:
        _add_chart_slide(prs, "Top Channels by Return Value", r["channels"], theme)

    try:
        insight = rag.query("Summarize last week return analysis insights for management")
    except Exception:
        insight = "AI Insight generation temporarily unavailable."

    _add_ai_summary_slide(prs, theme, insight)

    # Add Hive Mind context if provided
    if additional_context:
        try:
            slide = prs.slides.add_slide(prs.slide_layouts[5])
            slide.shapes.title.text = "Hive Mind Context Summary"
            tx = slide.shapes.add_textbox(Inches(0.8), Inches(1.8), Inches(8.5), Inches(4))
            body = tx.text_frame
            body.text = additional_context[:2000]
            body.paragraphs[0].font.size = Pt(16)
            body.paragraphs[0].font.color.rgb = RGBColor(*_hex_to_rgb(theme.get("text", "#000000")))
            _add_logo_and_watermark(slide, theme)
        except Exception as e:
            print(f"⚠️ Could not add context slide: {e}")

    out_path = os.path.join(OUTPUT_DIR, "PPT", f"Returns_Report_v3_{theme_key}_theme.pptx")
    os.makedirs(os.path.dirname(out_path), exist_ok=True)
    prs.save(out_path)
    return out_path
