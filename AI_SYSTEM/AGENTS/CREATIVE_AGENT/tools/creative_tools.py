# =========================================================
# creative_tools.py — Creative Agent Utility Toolkit (PPT/PDF/CSV)
# =========================================================
import os, json
from datetime import datetime, timedelta
import pandas as pd
import matplotlib.pyplot as plt

# PPT
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor as RGB
from pptx.enum.text import PP_ALIGN

# PDF (Advanced Creative Builder)
from io import BytesIO
try:
    from reportlab.lib import colors
    from reportlab.lib.pagesizes import A4
    from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer, Table, TableStyle, Image
    from reportlab.lib.styles import getSampleStyleSheet
    HAVE_PDF = True
except Exception:
    HAVE_PDF = False

# ---------------------------------------------------------
# ⚙️ PATHS
# ---------------------------------------------------------
BASE_PATH = r"C:\Users\Sandeep\Desktop\roborana_ai_system\RoboRana_AI_Data"
OUTPUT_DIR = os.path.join(BASE_PATH, "OUTPUTS")
os.makedirs(OUTPUT_DIR, exist_ok=True)
MEMORY_FILE = os.path.join(BASE_PATH, "AI_SYSTEM", "MEMORY", "creative_profile.json")
os.makedirs(os.path.dirname(MEMORY_FILE), exist_ok=True)

# ---------------------------------------------------------
# 🎨 PERSONALITY PROFILE
# ---------------------------------------------------------
DEFAULT_PROFILE = {
    "preferred_style": "modern_corporate",
    "palette": {
        "bg": "#FFFFFF",
        "primary": "#1F4E79",
        "accent": "#6FA8DC",
        "text": "#1A1A1A",
        "muted": "#999999"
    },
    "font": "Calibri",
    "report_preferences": {
        "chart_round": True,
        "title_case": True,
        "export_format": "pptx"
    },
    "themes": {
        "modern_corporate": {"primary": "#1F4E79", "accent": "#6FA8DC", "bg": "#FFFFFF"},
        "floral": {"primary": "#9B51E0", "accent": "#F299C1", "bg": "#FFF8FB"},
        "midnight": {"primary": "#0B132B", "accent": "#1C2541", "bg": "#0B132B"},
        "sunset": {"primary": "#D86C70", "accent": "#F2C14E", "bg": "#FFF3E0"}
    }
}

def load_profile():
    if not os.path.exists(MEMORY_FILE):
        save_profile(DEFAULT_PROFILE)
        return DEFAULT_PROFILE
    try:
        with open(MEMORY_FILE, "r", encoding="utf-8") as f:
            return json.load(f)
    except Exception:
        save_profile(DEFAULT_PROFILE)
        return DEFAULT_PROFILE

def save_profile(data):
    with open(MEMORY_FILE, "w", encoding="utf-8") as f:
        json.dump(data, f, indent=2, ensure_ascii=False)

def choose_theme(profile, theme_hint: str | None):
    theme_key = (theme_hint or profile.get("preferred_style") or "modern_corporate").strip().lower()
    themes = profile.get("themes", {})
    if theme_key in themes:
        palette = themes[theme_key]
    else:
        palette = themes.get("modern_corporate", DEFAULT_PROFILE["palette"])
    merged = {**DEFAULT_PROFILE["palette"], **palette}
    return theme_key, merged

# ---------------------------------------------------------
# 📊 DATA HELPERS
# ---------------------------------------------------------
def _find_col(df: pd.DataFrame, candidates: list[str]) -> str | None:
    cols = {c.lower(): c for c in df.columns}
    for key in candidates:
        for c in cols:
            if key in c:
                return cols[c]
    return None

def load_sales_master():
    path = os.path.join(BASE_PATH, "DATA", "SALES", "Master", "Sales_Master.csv")
    if not os.path.exists(path):
        raise FileNotFoundError(f"Sales master not found: {path}")
    return pd.read_csv(path, on_bad_lines="skip", low_memory=False, encoding="utf-8-sig")

def load_returns_master():
    path_upd = os.path.join(BASE_PATH, "DATA", "RETURNS", "Master", "Return_Master_Updated.csv")
    path_std = os.path.join(BASE_PATH, "DATA", "RETURNS", "Master", "Return_Master.csv")
    path = path_upd if os.path.exists(path_upd) else path_std
    if not os.path.exists(path):
        raise FileNotFoundError(f"Return master not found: {path}")
    return pd.read_csv(path, on_bad_lines="skip", low_memory=False, encoding="utf-8-sig")

# ---------------------------------------------------------
# 🧮 KPI SUMMARIES
# ---------------------------------------------------------
def summarize_sales(df: pd.DataFrame, days: int = 7):
    dt_col = _find_col(df, ["order date", "created", "placed on"])
    price_col = _find_col(df, ["selling price", "total price", "final amount"])
    channel_col = _find_col(df, ["channel name", "portal", "platform"])
    order_code_col = _find_col(df, ["display order code", "seller order id", "order code"])

    if dt_col:
        s = pd.to_datetime(df[dt_col], errors="coerce", dayfirst=False)
        df = df[s >= (datetime.now() - timedelta(days=days))]

    total_rows = len(df)
    total_value = pd.to_numeric(df[price_col], errors="coerce").fillna(0).sum() if price_col else 0.0
    orders = df[order_code_col].nunique() if order_code_col else total_rows
    aov = (total_value / orders) if orders else 0

    channel_summary = {}
    if channel_col:
        if price_col:
            channel_summary = df.groupby(channel_col)[price_col].sum().sort_values(ascending=False).head(10).to_dict()
        else:
            channel_summary = df.groupby(channel_col).size().sort_values(ascending=False).head(10).to_dict()

    return {
        "total_rows": total_rows,
        "total_orders": orders,
        "total_value": round(total_value, 2),
        "aov": round(aov, 2),
        "channels": channel_summary
    }

def summarize_returns(df: pd.DataFrame, days: int = 7):
    dt_col = _find_col(df, ["date", "created", "updated"])
    total_col = _find_col(df, ["total", "amount", "value", "selling price"])
    channel_col = _find_col(df, ["channel", "channel entry", "channel name"])
    qty_col = _find_col(df, ["qty", "quantity"])

    if dt_col:
        s = pd.to_datetime(df[dt_col], errors="coerce", dayfirst=True)
        df = df[s >= (datetime.now() - timedelta(days=days))]

    qty = int(pd.to_numeric(df[qty_col], errors="coerce").fillna(0).sum()) if qty_col else len(df)
    value = float(pd.to_numeric(df[total_col], errors="coerce").fillna(0).sum()) if total_col else 0.0

    channel_summary = {}
    if channel_col:
        if total_col:
            channel_summary = df.groupby(channel_col)[total_col].sum().sort_values(ascending=False).head(10).to_dict()
        else:
            channel_summary = df.groupby(channel_col).size().sort_values(ascending=False).head(10).to_dict()

    return {"rows": len(df), "qty": qty, "value": round(value, 2), "channels": channel_summary}

# ---------------------------------------------------------
# 🖼️ PPT BUILDERS (same as before)
# ---------------------------------------------------------
def _hex_to_rgb(hex_color):
    hex_color = hex_color.lstrip("#")
    return tuple(int(hex_color[i:i+2], 16) for i in (0, 2, 4))

def _title_slide(prs, title, subtitle, theme):
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = RGB(*_hex_to_rgb(theme.get("bg", "#FFFFFF")))
    slide.shapes.title.text = title
    slide.shapes.title.text_frame.paragraphs[0].font.size = Pt(40)
    slide.shapes.title.text_frame.paragraphs[0].font.bold = True
    slide.shapes.title.text_frame.paragraphs[0].font.color.rgb = RGB(*_hex_to_rgb(theme.get("primary", "#000000")))
    slide.placeholders[1].text = subtitle
    slide.placeholders[1].text_frame.paragraphs[0].font.size = Pt(18)
    slide.placeholders[1].text_frame.paragraphs[0].font.color.rgb = RGB(*_hex_to_rgb(theme.get("accent", "#777777")))

def _kpi_slide(prs, heading, kpis, theme=None):
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = RGB(*_hex_to_rgb(theme.get("bg", "#FFFFFF")))
    slide.shapes.title.text = heading
    tf = slide.shapes.title.text_frame
    tf.paragraphs[0].font.size = Pt(28)
    tf.paragraphs[0].font.color.rgb = RGB(*_hex_to_rgb(theme.get("primary", "#000000")))
    left, top, width, height = Inches(0.6), Inches(1.7), Inches(9.0), Inches(4.0)
    tx = slide.shapes.add_textbox(left, top, width, height).text_frame
    for label, value in kpis:
        p = tx.add_paragraph()
        p.text = f"{label}: {value}"
        p.font.size = Pt(20)
        if theme:
            p.font.color.rgb = RGB(*_hex_to_rgb(theme.get("text", "#000000")))

def _bar_chart(path_png, title, series):
    fig = plt.figure()
    labels, vals = list(series.keys())[:10], list(series.values())[:10]
    plt.barh(labels, vals)
    plt.title(title)
    plt.tight_layout()
    fig.savefig(path_png, dpi=200)
    plt.close(fig)

def _image_slide(prs, heading, img_path, theme=None):
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = RGB(*_hex_to_rgb(theme.get("bg", "#FFFFFF")))
    slide.shapes.title.text = heading
    slide.shapes.add_picture(img_path, Inches(0.6), Inches(1.6), width=Inches(9.0))

def build_sales_ppt(days=7, theme_hint=None):
    profile = load_profile()
    theme_key, theme = choose_theme(profile, theme_hint)
    df = load_sales_master()
    s = summarize_sales(df, days)
    prs = Presentation()
    now = datetime.now().strftime("%Y-%m-%d %H:%M")
    _title_slide(prs, f"Sales Performance — Last {days} Days", f"Theme: {theme_key} • Generated {now}", theme)
    kpis = [("Orders", f"{s['total_orders']:,}"), ("Rows", f"{s['total_rows']:,}"),
            ("Revenue", f"₹{s['total_value']:,.2f}"), ("AOV", f"₹{s['aov']:,.2f}")]
    _kpi_slide(prs, "Key Metrics", kpis, theme)
    tmp = os.path.join(OUTPUT_DIR, "TEMP"); os.makedirs(tmp, exist_ok=True)
    ch_png = os.path.join(tmp, "sales_channels.png")
    if s["channels"]:
        _bar_chart(ch_png, "Revenue by Channel", s["channels"])
        _image_slide(prs, "Top Channels", ch_png, theme)
    out_path = os.path.join(OUTPUT_DIR, "PPT", f"Sales_Report_last_{days}_days.pptx")
    os.makedirs(os.path.dirname(out_path), exist_ok=True)
    prs.save(out_path)
    return out_path

def build_returns_ppt(days=7, theme_hint=None):
    profile = load_profile()
    theme_key, theme = choose_theme(profile, theme_hint)
    df = load_returns_master()
    r = summarize_returns(df, days)
    prs = Presentation()
    now = datetime.now().strftime("%Y-%m-%d %H:%M")
    _title_slide(prs, f"Returns Summary — Last {days} Days", f"Theme: {theme_key} • Generated {now}", theme)
    kpis = [("Return Rows", f"{r['rows']:,}"), ("Return Qty", f"{r['qty']:,}"),
            ("Return Value", f"₹{r['value']:,.2f}")]
    _kpi_slide(prs, "Key Metrics", kpis, theme)
    tmp = os.path.join(OUTPUT_DIR, "TEMP"); os.makedirs(tmp, exist_ok=True)
    ch_png = os.path.join(tmp, "returns_channels.png")
    if r["channels"]:
        _bar_chart(ch_png, "Returns by Channel", r["channels"])
        _image_slide(prs, "Top Return Channels", ch_png, theme)
    out_path = os.path.join(OUTPUT_DIR, "PPT", f"Returns_Report_last_{days}_days.pptx")
    os.makedirs(os.path.dirname(out_path), exist_ok=True)
    prs.save(out_path)
    return out_path

# ---------------------------------------------------------
# 🌈 ADVANCED CREATIVE PDF BUILDER
# ---------------------------------------------------------
def build_creative_pdf(title: str, summary_data: dict, theme: dict) -> str:
    if not HAVE_PDF:
        raise RuntimeError("PDF generation requires 'reportlab'. Install it with: pip install reportlab")

    out_path = os.path.join(OUTPUT_DIR, "PDF", f"{title.replace(' ', '_')}.pdf")
    os.makedirs(os.path.dirname(out_path), exist_ok=True)

    doc = SimpleDocTemplate(out_path, pagesize=A4)
    styles = getSampleStyleSheet()
    elements = []

    # Title
    elements.append(Paragraph(f"<font color='{theme['primary']}' size=20><b>{title}</b></font>", styles["Title"]))
    elements.append(Spacer(1, 12))

    # Metrics Table
    metrics = [
        ["Metric", "Value"],
        ["Orders", f"{summary_data.get('orders', 0):,}"],
        ["Rows", f"{summary_data.get('rows', 0):,}"],
        ["Revenue", f"₹{summary_data.get('revenue', 0):,.2f}"],
        ["AOV", f"₹{summary_data.get('aov', 0):,.2f}"],
    ]
    table = Table(metrics, colWidths=[150, 200])
    table.setStyle(TableStyle([
        ("BACKGROUND", (0, 0), (-1, 0), colors.HexColor(theme["accent"])),
        ("TEXTCOLOR", (0, 0), (-1, 0), colors.white),
        ("BACKGROUND", (0, 1), (-1, -1), colors.HexColor("#F9F9F9")),
        ("GRID", (0, 0), (-1, -1), 0.25, colors.gray),
        ("FONTNAME", (0, 0), (-1, 0), "Helvetica-Bold"),
    ]))
    elements.append(table)
    elements.append(Spacer(1, 20))

    # Channels Chart
    channels = summary_data.get("channels", {})
    if channels:
        plt.figure(figsize=(5, 3))
        plt.barh(list(channels.keys())[:10], list(channels.values())[:10], color=theme["primary"])
        plt.title("Top Channels", color=theme["text"])
        plt.tight_layout()
        buf = BytesIO()
        plt.savefig(buf, format="PNG"); plt.close()
        buf.seek(0)
        elements.append(Image(buf, width=400, height=250))
        elements.append(Spacer(1, 20))

    elements.append(Paragraph(
        f"<font color='{theme['muted']}'>Generated on {datetime.now().strftime('%Y-%m-%d %H:%M')}</font>",
        styles["Normal"]
    ))

    doc.build(elements)
    return out_path

# ---------------------------------------------------------
# 💾 CSV EXPORT
# ---------------------------------------------------------
def export_csv(df: pd.DataFrame, name: str) -> str:
    out_path = os.path.join(OUTPUT_DIR, "CSV", f"{name}.csv")
    os.makedirs(os.path.dirname(out_path), exist_ok=True)
    df.to_csv(out_path, index=False, encoding="utf-8-sig")
    return out_path
